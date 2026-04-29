import xml.etree.ElementTree as ET
from pathlib import Path
from zipfile import ZIP_DEFLATED, ZipFile

import pytest
from PIL import Image
from pptx import Presentation

from autofacodex.contracts import SlideModel
from autofacodex.tools.pptx_generate import generate_pptx
from autofacodex.tools.pptx_inspect import inspect_pptx_editability


def _model(slide: dict) -> SlideModel:
    return SlideModel(slides=[slide])


def _rewrite_pptx_entries(path: Path, updates: dict[str, bytes]) -> None:
    with ZipFile(path) as source:
        entries = {
            item.filename: source.read(item.filename)
            for item in source.infolist()
            if not item.is_dir()
        }
    entries.update(updates)

    with ZipFile(path, "w", ZIP_DEFLATED) as target:
        for name, data in entries.items():
            target.writestr(name, data)


def _swap_presentation_slide_order(path: Path) -> None:
    xml_name = "ppt/presentation.xml"
    p_ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    r_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    ET.register_namespace("p", p_ns)
    ET.register_namespace("r", r_ns)

    with ZipFile(path) as source:
        presentation_xml = source.read(xml_name)
    root = ET.fromstring(presentation_xml)
    slide_id_list = root.find(f".//{{{p_ns}}}sldIdLst")
    assert slide_id_list is not None
    slide_ids = list(slide_id_list)
    assert len(slide_ids) == 2
    slide_id_list[:] = [slide_ids[1], slide_ids[0]]
    _rewrite_pptx_entries(
        path,
        {xml_name: ET.tostring(root, encoding="utf-8", xml_declaration=True)},
    )


def _add_orphan_slide(path: Path) -> None:
    rels_name = "ppt/_rels/presentation.xml.rels"
    content_types_name = "[Content_Types].xml"
    rels_ns = "http://schemas.openxmlformats.org/package/2006/relationships"
    content_types_ns = "http://schemas.openxmlformats.org/package/2006/content-types"
    ET.register_namespace("", rels_ns)

    with ZipFile(path) as source:
        slide_xml = source.read("ppt/slides/slide1.xml")
        rels_xml = source.read(rels_name)
        content_types_xml = source.read(content_types_name)

    rels_root = ET.fromstring(rels_xml)
    ET.SubElement(
        rels_root,
        f"{{{rels_ns}}}Relationship",
        {
            "Id": "rId999",
            "Type": "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide",
            "Target": "slides/slide99.xml",
        },
    )
    content_types_root = ET.fromstring(content_types_xml)
    ET.SubElement(
        content_types_root,
        f"{{{content_types_ns}}}Override",
        {
            "PartName": "/ppt/slides/slide99.xml",
            "ContentType": "application/vnd.openxmlformats-officedocument.presentationml.slide+xml",
        },
    )
    _rewrite_pptx_entries(
        path,
        {
            "ppt/slides/slide99.xml": slide_xml,
            rels_name: ET.tostring(rels_root, encoding="utf-8", xml_declaration=True),
            content_types_name: ET.tostring(
                content_types_root, encoding="utf-8", xml_declaration=True
            ),
        },
    )


def test_inspect_pptx_returns_editable_text_content(tmp_path: Path):
    model = _model(
        {
            "page_number": 1,
            "size": {"width": 10, "height": 7.5},
            "elements": [
                {
                    "id": "title",
                    "type": "text",
                    "text": "Editable Title",
                    "x": 1,
                    "y": 1,
                    "w": 4,
                    "h": 0.6,
                    "style": {"font_size": 24},
                }
            ],
            "raster_fallback_regions": [],
        }
    )
    output = tmp_path / "candidate.pptx"
    generate_pptx(model, output)

    inspection = inspect_pptx_editability(output)

    page = inspection["pages"][0]
    assert page["text"] == "Editable Title"
    assert page["text_runs"] == 1
    assert page["pictures"] == 0
    assert page["largest_picture_area_ratio"] == 0
    assert page["has_full_page_picture"] is False


def test_inspect_pptx_separates_text_runs_from_different_shapes(tmp_path: Path):
    model = _model(
        {
            "page_number": 1,
            "size": {"width": 10, "height": 7.5},
            "elements": [
                {
                    "id": "footer",
                    "type": "text",
                    "text": "Materials",
                    "x": 1,
                    "y": 1,
                    "w": 3,
                    "h": 0.4,
                },
                {
                    "id": "date",
                    "type": "text",
                    "text": "2024",
                    "x": 1,
                    "y": 2,
                    "w": 3,
                    "h": 0.4,
                },
            ],
            "raster_fallback_regions": [],
        }
    )
    output = tmp_path / "candidate.pptx"
    generate_pptx(model, output)

    inspection = inspect_pptx_editability(output)

    page = inspection["pages"][0]
    assert page["text"] == "Materials\n2024"
    assert page["text_runs"] == 2


def test_inspect_pptx_detects_full_page_picture(tmp_path: Path):
    image_path = tmp_path / "page.png"
    Image.new("RGB", (1200, 675), color=(240, 240, 240)).save(image_path)
    model = _model(
        {
            "page_number": 1,
            "size": {"width": 13.333, "height": 7.5},
            "elements": [
                {
                    "id": "page-image",
                    "type": "image",
                    "source": str(image_path),
                    "x": 0,
                    "y": 0,
                    "w": 13.333,
                    "h": 7.5,
                }
            ],
            "raster_fallback_regions": [],
        }
    )
    output = tmp_path / "candidate.pptx"
    generate_pptx(model, output)

    inspection = inspect_pptx_editability(output)

    page = inspection["pages"][0]
    assert page["pictures"] == 1
    assert page["largest_picture_area_ratio"] > 0.98
    assert page["has_full_page_picture"] is True


def test_inspect_pptx_detects_tiled_full_page_pictures(tmp_path: Path):
    image_path = tmp_path / "tile.png"
    Image.new("RGB", (600, 400), color=(240, 240, 240)).save(image_path)
    elements = []
    for index, (x, y) in enumerate(
        [(0, 0), (5, 0), (0, 3.75), (5, 3.75)], start=1
    ):
        elements.append(
            {
                "id": f"tile-{index}",
                "type": "image",
                "source": str(image_path),
                "x": x,
                "y": y,
                "w": 5,
                "h": 3.75,
            }
        )
    model = _model(
        {
            "page_number": 1,
            "size": {"width": 10, "height": 7.5},
            "elements": elements,
            "raster_fallback_regions": [],
        }
    )
    output = tmp_path / "candidate.pptx"
    generate_pptx(model, output)

    inspection = inspect_pptx_editability(output)

    page = inspection["pages"][0]
    assert page["pictures"] == 4
    assert page["largest_picture_area_ratio"] == pytest.approx(0.25, abs=0.01)
    assert page["picture_coverage_ratio"] > 0.98
    assert page["has_full_page_picture"] is True


def test_inspect_pptx_reports_text_box_count_and_geometries(tmp_path: Path):
    model = _model(
        {
            "page_number": 1,
            "size": {"width": 10, "height": 7.5},
            "elements": [
                {
                    "id": "title",
                    "type": "text",
                    "text": "Editable Title",
                    "x": 1,
                    "y": 1,
                    "w": 4,
                    "h": 0.6,
                    "style": {"font_size": 24},
                },
                {
                    "id": "subtitle",
                    "type": "text",
                    "text": "Editable Subtitle",
                    "x": 2,
                    "y": 2,
                    "w": 3,
                    "h": 0.4,
                    "style": {"font_size": 14},
                },
            ],
            "raster_fallback_regions": [],
        }
    )
    output = tmp_path / "candidate.pptx"
    generate_pptx(model, output)

    inspection = inspect_pptx_editability(output)

    page = inspection["pages"][0]
    assert page["text_box_count"] == 2
    assert page["text_box_geometries"] == [
        {"x": pytest.approx(1), "y": pytest.approx(1), "w": pytest.approx(4), "h": pytest.approx(0.6)},
        {"x": pytest.approx(2), "y": pytest.approx(2), "w": pytest.approx(3), "h": pytest.approx(0.4)},
    ]


def test_inspect_pptx_uses_presentation_slide_order(tmp_path: Path):
    output = tmp_path / "swapped-order.pptx"
    presentation = Presentation()
    blank = presentation.slide_layouts[6]
    slide_one = presentation.slides.add_slide(blank)
    slide_one.shapes.add_textbox(914400, 914400, 914400, 914400).text = "First XML"
    slide_two = presentation.slides.add_slide(blank)
    for index in range(20):
        slide_two.shapes.add_textbox(
            914400,
            914400 + index * 10000,
            914400,
            914400,
        ).text = f"Second XML {index}"
    presentation.save(output)
    _swap_presentation_slide_order(output)

    display_shape_counts = [len(slide.shapes) for slide in Presentation(output).slides]
    inspection = inspect_pptx_editability(output)

    assert display_shape_counts == [20, 1]
    assert [page["slide"] for page in inspection["pages"]] == [
        "ppt/slides/slide2.xml",
        "ppt/slides/slide1.xml",
    ]
    assert [page["shapes"] for page in inspection["pages"]] == display_shape_counts


def test_inspect_pptx_uses_numeric_fallback_when_presentation_xml_is_malformed(
    tmp_path: Path,
):
    output = tmp_path / "malformed-presentation.pptx"
    presentation = Presentation()
    blank = presentation.slide_layouts[6]
    slide = presentation.slides.add_slide(blank)
    slide.shapes.add_textbox(914400, 914400, 914400, 914400).text = "Still readable"
    presentation.save(output)
    _rewrite_pptx_entries(output, {"ppt/presentation.xml": b"<presentation>"})

    inspection = inspect_pptx_editability(output)

    assert [page["slide"] for page in inspection["pages"]] == ["ppt/slides/slide1.xml"]
    page = inspection["pages"][0]
    assert page["size"] == {"width": 10.0, "height": 7.5}
    assert page["text"] == "Still readable"


def test_inspect_pptx_ignores_orphan_slide_xml_when_display_order_resolves(
    tmp_path: Path,
):
    output = tmp_path / "orphan-slide.pptx"
    presentation = Presentation()
    blank = presentation.slide_layouts[6]
    slide = presentation.slides.add_slide(blank)
    slide.shapes.add_textbox(914400, 914400, 914400, 914400).text = "Displayed"
    presentation.save(output)
    _add_orphan_slide(output)

    inspection = inspect_pptx_editability(output)

    assert [page["slide"] for page in inspection["pages"]] == ["ppt/slides/slide1.xml"]


def test_inspect_pptx_reports_slide_name_for_malformed_slide_xml(tmp_path: Path):
    output = tmp_path / "malformed-slide.pptx"
    presentation = Presentation()
    blank = presentation.slide_layouts[6]
    presentation.slides.add_slide(blank)
    presentation.save(output)
    _rewrite_pptx_entries(output, {"ppt/slides/slide1.xml": b"<slide>"})

    with pytest.raises(RuntimeError, match="ppt/slides/slide1.xml"):
        inspect_pptx_editability(output)
