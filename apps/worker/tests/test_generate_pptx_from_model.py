from pathlib import Path

from PIL import Image
from pptx import Presentation

from autofacodex.contracts import SlideModel
from autofacodex.tools.generate_pptx_from_model import generate_from_model


def test_generate_from_model_uses_task_root_as_default_asset_root(tmp_path: Path):
    task_dir = tmp_path / "task"
    model_path = task_dir / "slides" / "slide-model.v2.json"
    output_path = task_dir / "output" / "candidate.v2.pptx"
    image_path = task_dir / "extracted" / "objects" / "images" / "logo.png"
    image_path.parent.mkdir(parents=True)
    model_path.parent.mkdir(parents=True)
    Image.new("RGB", (10, 10), color=(220, 20, 60)).save(image_path)
    model = SlideModel(
        slides=[
            {
                "page_number": 1,
                "size": {"width": 10, "height": 7.5},
                "elements": [
                    {
                        "id": "image-1",
                        "type": "image",
                        "source": "extracted/objects/images/logo.png",
                        "x": 1,
                        "y": 1,
                        "w": 2,
                        "h": 2,
                    }
                ],
                "raster_fallback_regions": [],
            }
        ]
    )
    model_path.write_text(model.model_dump_json(indent=2), encoding="utf-8")

    result = generate_from_model(model_path, output_path)

    assert result == output_path
    presentation = Presentation(output_path)
    assert len(presentation.slides) == 1
    assert len(presentation.slides[0].shapes) == 1
