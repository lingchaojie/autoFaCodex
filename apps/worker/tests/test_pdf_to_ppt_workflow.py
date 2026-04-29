import subprocess
from pathlib import Path

import fitz
import pytest
from pptx import Presentation
from reportlab.pdfgen import canvas

from autofacodex.config import WorkerConfig
from autofacodex.contracts import SlideModel, ValidatorReport
import autofacodex.workflows.pdf_to_ppt as workflow
from autofacodex.workflows.pdf_to_ppt import run_pdf_to_ppt


def make_pdf(path: Path) -> None:
    c = canvas.Canvas(str(path), pagesize=(400, 300))
    c.setFont("Helvetica", 24)
    c.drawString(72, 220, "Workflow Title 1")
    c.showPage()
    c.setFont("Helvetica", 24)
    c.drawString(72, 220, "Workflow Title 2")
    c.save()


def _write_validator_report(task_dir: Path, attempt: int, page_count: int) -> ValidatorReport:
    report = ValidatorReport(
        task_id=task_dir.name,
        attempt=attempt,
        aggregate_status="pass",
        pages=[
            {
                "page_number": page_number,
                "status": "pass",
                "visual_score": 1.0,
                "editable_score": 1.0,
                "text_coverage_score": 1.0,
                "raster_fallback_ratio": 0,
                "issues": [],
            }
            for page_number in range(1, page_count + 1)
        ],
    )
    (task_dir / "reports").mkdir(parents=True, exist_ok=True)
    (task_dir / "reports" / f"validator.v{attempt}.json").write_text(
        report.model_dump_json(indent=2),
        encoding="utf-8",
    )
    return report


def _write_next_slide_model_text(
    task_dir: Path, source_attempt: int, target_attempt: int, text: str
) -> None:
    source_model = SlideModel.model_validate_json(
        (task_dir / "slides" / f"slide-model.v{source_attempt}.json").read_text(
            encoding="utf-8"
        )
    )
    data = source_model.model_dump()
    for slide in data["slides"]:
        for element in slide["elements"]:
            if element["type"] == "text":
                element["text"] = text
                target_model = SlideModel.model_validate(data)
                (task_dir / "slides" / f"slide-model.v{target_attempt}.json").write_text(
                    target_model.model_dump_json(indent=2),
                    encoding="utf-8",
                )
                return
    raise AssertionError("expected the generated slide model to contain text")


def _write_minimal_slide_model(task_dir: Path, attempt: int = 1) -> None:
    model = SlideModel(
        slides=[
            {
                "page_number": 1,
                "size": {"width": 10, "height": 7.5},
                "elements": [
                    {
                        "id": "title",
                        "type": "text",
                        "text": "Title",
                        "x": 1,
                        "y": 1,
                        "w": 4,
                        "h": 0.5,
                    }
                ],
                "raster_fallback_regions": [],
            }
        ]
    )
    (task_dir / "slides").mkdir(parents=True, exist_ok=True)
    (task_dir / "slides" / f"slide-model.v{attempt}.json").write_text(
        model.model_dump_json(indent=2),
        encoding="utf-8",
    )


def test_run_pdf_to_ppt_creates_candidate_report_and_slide_model(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
):
    task_dir = tmp_path / "task_1"
    task_dir.mkdir()
    make_pdf(task_dir / "input.pdf")

    def fake_validate_candidate(received_task_dir: Path, attempt: int):
        assert received_task_dir == task_dir
        assert attempt == 1
        return _write_validator_report(received_task_dir, attempt, page_count=2)

    monkeypatch.setattr(
        workflow, "validate_candidate", fake_validate_candidate, raising=False
    )

    run_pdf_to_ppt(task_dir)

    assert (task_dir / "output" / "candidate.v1.pptx").is_file()
    assert (task_dir / "output" / "final.pptx").is_file()
    assert (task_dir / "reports" / "validator.v1.json").is_file()
    assert not (task_dir / "slides" / "slide-model.high-fidelity.v1.json").exists()
    assert (task_dir / "slides" / "slide-model.v1.json").is_file()
    assert (task_dir / "slides" / "slide-model.final.v1.json").is_file()

    model = SlideModel.model_validate_json(
        (task_dir / "slides" / "slide-model.v1.json").read_text(encoding="utf-8")
    )
    report = ValidatorReport.model_validate_json(
        (task_dir / "reports" / "validator.v1.json").read_text(encoding="utf-8")
    )
    presentation = Presentation(task_dir / "output" / "candidate.v1.pptx")
    final_presentation = Presentation(task_dir / "output" / "final.pptx")
    final_model = SlideModel.model_validate_json(
        (task_dir / "slides" / "slide-model.final.v1.json").read_text(encoding="utf-8")
    )

    assert [slide.page_number for slide in model.slides] == [1, 2]
    assert final_model == model
    assert [page.page_number for page in report.pages] == [1, 2]
    assert len(presentation.slides) == len(model.slides) == len(report.pages)
    assert len(final_presentation.slides) == len(model.slides)


def test_run_pdf_to_ppt_promotes_accepted_semantic_table_repair_to_final(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
):
    task_dir = tmp_path / "task_1"
    task_dir.mkdir()
    make_pdf(task_dir / "input.pdf")

    def fake_validate_candidate(received_task_dir: Path, attempt: int):
        return _write_validator_report(received_task_dir, attempt, page_count=2)

    def fake_semantic_table_repair(received_task_dir: Path, *, source_attempt: int):
        assert received_task_dir == task_dir
        assert source_attempt == 1
        _write_next_slide_model_text(task_dir, source_attempt, 2, "Accepted Semantic V2")
        return {"status": "accepted", "source_attempt": source_attempt, "target_attempt": 2}

    monkeypatch.setattr(
        workflow, "validate_candidate", fake_validate_candidate, raising=False
    )
    monkeypatch.setattr(
        workflow,
        "upgrade_semantic_tables_with_guard",
        fake_semantic_table_repair,
        raising=False,
    )

    run_pdf_to_ppt(task_dir)

    assert (task_dir / "slides" / "slide-model.final.v2.json").is_file()
    assert not (task_dir / "slides" / "slide-model.final.v1.json").exists()
    final_model = SlideModel.model_validate_json(
        (task_dir / "slides" / "slide-model.final.v2.json").read_text(encoding="utf-8")
    )
    assert any(
        element.text == "Accepted Semantic V2"
        for slide in final_model.slides
        for element in slide.elements
    )


def test_run_pdf_to_ppt_keeps_source_final_when_semantic_table_repair_rejects(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
):
    task_dir = tmp_path / "task_1"
    task_dir.mkdir()
    make_pdf(task_dir / "input.pdf")

    def fake_validate_candidate(received_task_dir: Path, attempt: int):
        return _write_validator_report(received_task_dir, attempt, page_count=2)

    def fake_semantic_table_repair(received_task_dir: Path, *, source_attempt: int):
        assert received_task_dir == task_dir
        _write_next_slide_model_text(task_dir, source_attempt, 2, "Rejected Semantic V2")
        return {"status": "rejected", "source_attempt": source_attempt, "target_attempt": 2}

    monkeypatch.setattr(
        workflow, "validate_candidate", fake_validate_candidate, raising=False
    )
    monkeypatch.setattr(
        workflow,
        "upgrade_semantic_tables_with_guard",
        fake_semantic_table_repair,
        raising=False,
    )

    run_pdf_to_ppt(task_dir)

    assert (task_dir / "slides" / "slide-model.v2.json").is_file()
    assert (task_dir / "slides" / "slide-model.final.v1.json").is_file()
    assert not (task_dir / "slides" / "slide-model.final.v2.json").exists()
    final_model = SlideModel.model_validate_json(
        (task_dir / "slides" / "slide-model.final.v1.json").read_text(encoding="utf-8")
    )
    assert all(
        element.text != "Rejected Semantic V2"
        for slide in final_model.slides
        for element in slide.elements
    )


def test_run_pdf_to_ppt_initial_uses_real_candidate_validation(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
):
    task_dir = tmp_path / "task_1"
    task_dir.mkdir()
    make_pdf(task_dir / "input.pdf")
    calls = []

    def fake_validate_candidate(received_task_dir: Path, attempt: int):
        calls.append((received_task_dir, attempt))
        report = ValidatorReport(
            task_id=received_task_dir.name,
            attempt=attempt,
            aggregate_status="repair_needed",
            pages=[
                {
                    "page_number": 1,
                    "status": "repair_needed",
                    "visual_score": 0.5,
                    "editable_score": 1.0,
                    "text_coverage_score": 1.0,
                    "raster_fallback_ratio": 0,
                    "issues": [],
                }
            ],
        )
        (received_task_dir / "reports").mkdir(parents=True, exist_ok=True)
        (received_task_dir / "reports" / "validator.v1.json").write_text(
            report.model_dump_json(indent=2),
            encoding="utf-8",
        )
        return report

    monkeypatch.setattr(
        workflow, "validate_candidate", fake_validate_candidate, raising=False
    )

    run_pdf_to_ppt(task_dir)

    assert calls == [(task_dir, 1)]
    assert not (task_dir / "slides" / "slide-model.high-fidelity.v1.json").exists()
    final_model = SlideModel.model_validate_json(
        (task_dir / "slides" / "slide-model.final.v1.json").read_text(encoding="utf-8")
    )
    assert all(
        element.style.get("role") != "visual_background"
        for slide in final_model.slides
        for element in slide.elements
    )
    assert all(not slide.raster_fallback_regions for slide in final_model.slides)


def test_run_pdf_to_ppt_initial_runs_ai_repair_when_manifest_allows(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
):
    task_dir = tmp_path / "task_1"
    task_dir.mkdir()
    make_pdf(task_dir / "input.pdf")
    (task_dir / "task-manifest.json").write_text(
        """
        {
          "task_id": "task_1",
          "workflow_type": "pdf_to_ppt",
          "input_pdf": "input.pdf",
          "attempt": 1,
          "max_attempts": 2
        }
        """,
        encoding="utf-8",
    )
    repair_calls = []

    def fake_validate_candidate(received_task_dir: Path, attempt: int):
        report = ValidatorReport(
            task_id=received_task_dir.name,
            attempt=attempt,
            aggregate_status="repair_needed",
            pages=[
                {
                    "page_number": 1,
                    "status": "repair_needed",
                    "visual_score": 0.5,
                    "editable_score": 1.0,
                    "text_coverage_score": 1.0,
                    "raster_fallback_ratio": 0,
                    "issues": [],
                }
            ],
        )
        (received_task_dir / "reports").mkdir(parents=True, exist_ok=True)
        (received_task_dir / "reports" / f"validator.v{attempt}.json").write_text(
            report.model_dump_json(indent=2),
            encoding="utf-8",
        )
        return report

    def fake_run_repair(received_task_dir: Path):
        repair_calls.append(received_task_dir)
        _write_validator_report(received_task_dir, attempt=2, page_count=1)

    monkeypatch.setattr(
        workflow, "validate_candidate", fake_validate_candidate, raising=False
    )
    monkeypatch.setattr(workflow, "_run_repair", fake_run_repair, raising=False)

    run_pdf_to_ppt(task_dir)

    assert repair_calls == [task_dir]


def test_run_pdf_to_ppt_initial_stops_repair_at_manifest_attempt_limit(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
):
    task_dir = tmp_path / "task_1"
    task_dir.mkdir()
    make_pdf(task_dir / "input.pdf")
    (task_dir / "task-manifest.json").write_text(
        """
        {
          "task_id": "task_1",
          "workflow_type": "pdf_to_ppt",
          "input_pdf": "input.pdf",
          "attempt": 1,
          "max_attempts": 1
        }
        """,
        encoding="utf-8",
    )

    def fake_validate_candidate(received_task_dir: Path, attempt: int):
        report = ValidatorReport(
            task_id=received_task_dir.name,
            attempt=attempt,
            aggregate_status="repair_needed",
            pages=[
                {
                    "page_number": 1,
                    "status": "repair_needed",
                    "visual_score": 0.5,
                    "editable_score": 1.0,
                    "text_coverage_score": 1.0,
                    "raster_fallback_ratio": 0,
                    "issues": [],
                }
            ],
        )
        (received_task_dir / "reports").mkdir(parents=True, exist_ok=True)
        (received_task_dir / "reports" / f"validator.v{attempt}.json").write_text(
            report.model_dump_json(indent=2),
            encoding="utf-8",
        )
        return report

    def fail_run_repair(_received_task_dir: Path):
        raise AssertionError("repair should not run after max_attempts is reached")

    monkeypatch.setattr(
        workflow, "validate_candidate", fake_validate_candidate, raising=False
    )
    monkeypatch.setattr(workflow, "_run_repair", fail_run_repair, raising=False)

    run_pdf_to_ppt(task_dir)

    report = ValidatorReport.model_validate_json(
        (task_dir / "reports" / "validator.v1.json").read_text(encoding="utf-8")
    )
    assert report.aggregate_status == "repair_needed"


def test_run_pdf_to_ppt_rejects_invalid_pdf(tmp_path: Path):
    task_dir = tmp_path / "task_1"
    task_dir.mkdir()
    (task_dir / "input.pdf").write_bytes(b"%PDF-1.4")

    with pytest.raises(fitz.FileDataError):
        run_pdf_to_ppt(task_dir)

    assert not (task_dir / "output" / "candidate.v1.pptx").exists()


def test_run_pdf_to_ppt_rejects_missing_pdf_renders(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
):
    task_dir = tmp_path / "task_1"
    task_dir.mkdir()
    make_pdf(task_dir / "input.pdf")

    monkeypatch.setattr("autofacodex.workflows.pdf_to_ppt.render_pdf_pages", lambda *_: [])

    with pytest.raises(RuntimeError, match="Expected 2 PDF renders"):
        run_pdf_to_ppt(task_dir)

    assert not (task_dir / "output" / "candidate.v1.pptx").exists()
    assert not (task_dir / "reports" / "validator.v1.json").exists()


def test_run_pdf_to_ppt_rejects_nonexistent_pdf_render_paths(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
):
    task_dir = tmp_path / "task_1"
    task_dir.mkdir()
    make_pdf(task_dir / "input.pdf")

    monkeypatch.setattr(
        "autofacodex.workflows.pdf_to_ppt.render_pdf_pages",
        lambda *_: [task_dir / "missing-1.png", task_dir / "missing-2.png"],
    )

    with pytest.raises(RuntimeError, match="PDF render does not exist"):
        run_pdf_to_ppt(task_dir)

    assert not (task_dir / "output" / "candidate.v1.pptx").exists()
    assert not (task_dir / "reports" / "validator.v1.json").exists()


def test_run_pdf_to_ppt_repair_invokes_runner_then_validator_and_logs_results(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
):
    task_dir = tmp_path / "task_1"
    task_dir.mkdir()
    _write_minimal_slide_model(task_dir, attempt=1)
    _write_validator_report(task_dir, attempt=1, page_count=1)
    (task_dir / "task-manifest.json").write_text("{}", encoding="utf-8")
    (task_dir / "conversation").mkdir()
    (task_dir / "conversation" / "messages.jsonl").write_text(
        '{"role":"user","content":"Fix slide 2","createdAt":"2026-04-26T12:00:00.000Z"}\n',
        encoding="utf-8",
    )
    config = WorkerConfig(
        redis_url="redis://example",
        shared_tasks_dir=tmp_path,
        codex_home=tmp_path / "codex-home",
        codex_bin="codex-test",
    )
    calls = []
    results = [
        subprocess.CompletedProcess(
            args=["codex-test"], returncode=0, stdout="runner stdout", stderr="runner stderr"
        ),
        subprocess.CompletedProcess(
            args=["codex-test"],
            returncode=0,
            stdout="validator stdout",
            stderr="validator stderr",
        ),
    ]

    def fail_initial_extraction(*_args, **_kwargs):
        raise AssertionError("repair mode must not run deterministic initial extraction")

    def fake_run_codex_agent(invocation, message: str):
        calls.append((invocation, message))
        if invocation.role == "runner":
            _write_minimal_slide_model(task_dir, attempt=2)
            (task_dir / "output").mkdir(parents=True, exist_ok=True)
            (task_dir / "output" / "candidate.v2.pptx").write_bytes(b"pptx")
            (task_dir / "reports" / "runner-repair.v2.json").write_text(
                "{}",
                encoding="utf-8",
            )
        if invocation.role == "validator":
            _write_validator_report(task_dir, attempt=2, page_count=1)
        return results.pop(0)

    monkeypatch.setattr(workflow, "extract_pdf", fail_initial_extraction)
    monkeypatch.setattr(workflow, "load_config", lambda: config, raising=False)
    monkeypatch.setattr(workflow, "run_codex_agent", fake_run_codex_agent, raising=False)

    run_pdf_to_ppt(task_dir, mode="repair")

    assert [call[0].role for call in calls] == ["runner", "validator"]
    assert [call[0].task_dir for call in calls] == [task_dir, task_dir]
    assert calls[0][0].codex_home == config.codex_home
    assert calls[0][0].codex_bin == "codex-test"
    assert calls[0][0].system_prompt.as_posix().endswith(
        "agent_assets/runner/runner.system.md"
    )
    assert calls[0][0].skill_dir.as_posix().endswith("agent_assets/runner")
    assert calls[1][0].system_prompt.as_posix().endswith(
        "agent_assets/validator/validator.system.md"
    )
    assert calls[1][0].skill_dir.as_posix().endswith("agent_assets/validator")
    assert "task-manifest.json" in calls[0][1]
    assert "latest reports/validator.v*.json" in calls[0][1]
    assert "latest slides/slide-model.v*.json" in calls[0][1]
    assert "conversation/messages.jsonl" in calls[0][1]
    assert "initial PDF extraction or generation" not in calls[0][1]
    assert "produce a revised slide model" in calls[0][1]
    assert "single bounded repair pass" in calls[0][1]
    assert "Do not render, diff, or visually validate" in calls[0][1]
    assert "PYTHONPATH" in calls[0][1]
    assert "autofacodex.tools.generate_pptx_from_model" in calls[0][1]
    assert "autofacodex.tools.semantic_table_repair" in calls[0][1]
    assert "validate every page after repair" in calls[1][1]
    assert "autofacodex.tools.pptx_inspect" in calls[1][1]
    assert "autofacodex.tools.visual_diff" in calls[1][1]
    assert "reports/validator.vN.json" in calls[1][1]
    assert "candidate.v2.pptx" in calls[0][1]
    assert "validator.v2.json" in calls[1][1]
    assert (task_dir / "logs" / "runner-repair.log").read_text(encoding="utf-8") == (
        "returncode: 0\nstdout:\nrunner stdout\nstderr:\nrunner stderr\n"
    )
    assert (task_dir / "logs" / "validator-repair.log").read_text(encoding="utf-8") == (
        "returncode: 0\nstdout:\nvalidator stdout\nstderr:\nvalidator stderr\n"
    )


def test_run_pdf_to_ppt_repair_logs_runner_timeout(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
):
    task_dir = tmp_path / "task_1"
    task_dir.mkdir()
    _write_minimal_slide_model(task_dir, attempt=1)
    _write_validator_report(task_dir, attempt=1, page_count=1)
    config = WorkerConfig(
        redis_url="redis://example",
        shared_tasks_dir=tmp_path,
        codex_home=tmp_path / "codex-home",
        codex_bin="codex-test",
    )

    def fake_run_codex_agent(_invocation, _message: str):
        raise subprocess.TimeoutExpired(
            cmd=["codex-test"],
            timeout=900,
            output=b"partial stdout",
            stderr=b"partial stderr",
        )

    def fake_validate_candidate(received_task_dir: Path, attempt: int):
        return _write_validator_report(received_task_dir, attempt=attempt, page_count=1)

    monkeypatch.setattr(workflow, "load_config", lambda: config, raising=False)
    monkeypatch.setattr(workflow, "run_codex_agent", fake_run_codex_agent, raising=False)
    monkeypatch.setattr(
        workflow, "validate_candidate", fake_validate_candidate, raising=False
    )

    run_pdf_to_ppt(task_dir, mode="repair")

    runner_log = (task_dir / "logs" / "runner-repair.log").read_text(encoding="utf-8")
    validator_log = (task_dir / "logs" / "validator-repair.log").read_text(encoding="utf-8")
    assert "timeout_after_seconds: 900" in runner_log
    assert "stdout:\npartial stdout" in runner_log
    assert "stderr:\npartial stderr" in runner_log
    assert "deterministic runner fallback" in runner_log
    assert "timeout_after_seconds: 900" in validator_log
    assert "deterministic validator fallback" in validator_log
    assert (task_dir / "reports" / "validator.v2.json").is_file()


def test_run_pdf_to_ppt_repair_falls_back_when_agents_do_not_write_artifacts(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
):
    task_dir = tmp_path / "task_1"
    task_dir.mkdir()
    _write_minimal_slide_model(task_dir, attempt=1)
    _write_validator_report(task_dir, attempt=1, page_count=1)
    config = WorkerConfig(
        redis_url="redis://example",
        shared_tasks_dir=tmp_path,
        codex_home=tmp_path / "codex-home",
        codex_bin="codex-test",
    )
    calls = []
    fallback_calls = []
    validation_calls = []

    def fake_run_codex_agent(invocation, message: str):
        calls.append((invocation.role, message))
        if invocation.role == "runner":
            raise subprocess.TimeoutExpired(
                cmd=["codex-test"],
                timeout=900,
                output=b"runner partial stdout",
                stderr=b"runner partial stderr",
            )
        return subprocess.CompletedProcess(
            args=["codex-test"],
            returncode=0,
            stdout="validator stdout",
            stderr="validator skipped report",
        )

    def fake_deterministic_repair(
        received_task_dir: Path,
        *,
        source_attempt: int,
        target_attempt: int,
        reason: str,
    ):
        fallback_calls.append(
            (received_task_dir, source_attempt, target_attempt, reason)
        )
        _write_minimal_slide_model(received_task_dir, attempt=target_attempt)
        (received_task_dir / "output").mkdir(parents=True, exist_ok=True)
        (received_task_dir / "output" / "candidate.v2.pptx").write_bytes(b"pptx")
        (received_task_dir / "reports" / "runner-repair.v2.json").write_text(
            '{"mode":"deterministic_fallback"}',
            encoding="utf-8",
        )
        return {"mode": "deterministic_fallback"}

    def fake_validate_candidate(received_task_dir: Path, attempt: int):
        validation_calls.append((received_task_dir, attempt))
        return _write_validator_report(received_task_dir, attempt=attempt, page_count=1)

    monkeypatch.setattr(workflow, "load_config", lambda: config, raising=False)
    monkeypatch.setattr(workflow, "run_codex_agent", fake_run_codex_agent, raising=False)
    monkeypatch.setattr(
        workflow,
        "run_deterministic_runner_repair",
        fake_deterministic_repair,
        raising=False,
    )
    monkeypatch.setattr(
        workflow, "validate_candidate", fake_validate_candidate, raising=False
    )

    run_pdf_to_ppt(task_dir, mode="repair")

    assert [role for role, _message in calls] == ["runner", "validator"]
    assert fallback_calls == [
        (task_dir, 1, 2, "runner_timeout: Runner repair timed out after 900 seconds")
    ]
    assert validation_calls == [(task_dir, 2)]
    assert (task_dir / "slides" / "slide-model.v2.json").is_file()
    assert (task_dir / "output" / "candidate.v2.pptx").is_file()
    assert (task_dir / "reports" / "runner-repair.v2.json").is_file()
    assert (task_dir / "reports" / "validator.v2.json").is_file()
    assert "deterministic runner fallback" in (
        task_dir / "logs" / "runner-repair.log"
    ).read_text(encoding="utf-8")
    assert "deterministic validator fallback" in (
        task_dir / "logs" / "validator-repair.log"
    ).read_text(encoding="utf-8")


def test_run_pdf_to_ppt_repair_runner_failure_prevents_validator(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
):
    task_dir = tmp_path / "task_1"
    task_dir.mkdir()
    _write_minimal_slide_model(task_dir, attempt=1)
    _write_validator_report(task_dir, attempt=1, page_count=1)
    config = WorkerConfig(
        redis_url="redis://example",
        shared_tasks_dir=tmp_path,
        codex_home=tmp_path / "codex-home",
        codex_bin="codex-test",
    )
    calls = []

    def fake_run_codex_agent(invocation, message: str):
        calls.append((invocation, message))
        return subprocess.CompletedProcess(
            args=["codex-test"], returncode=7, stdout="runner stdout", stderr="runner failed"
        )

    monkeypatch.setattr(workflow, "load_config", lambda: config, raising=False)
    monkeypatch.setattr(workflow, "run_codex_agent", fake_run_codex_agent, raising=False)

    with pytest.raises(RuntimeError, match="Runner repair failed with return code 7"):
        run_pdf_to_ppt(task_dir, mode="repair")

    assert [call[0].role for call in calls] == ["runner"]
    assert (task_dir / "logs" / "runner-repair.log").read_text(encoding="utf-8") == (
        "returncode: 7\nstdout:\nrunner stdout\nstderr:\nrunner failed\n"
    )
    assert not (task_dir / "logs" / "validator-repair.log").exists()


def test_run_pdf_to_ppt_repair_falls_back_after_runner_success_without_required_artifacts(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
):
    task_dir = tmp_path / "task_1"
    task_dir.mkdir()
    _write_minimal_slide_model(task_dir, attempt=1)
    _write_validator_report(task_dir, attempt=1, page_count=1)
    config = WorkerConfig(
        redis_url="redis://example",
        shared_tasks_dir=tmp_path,
        codex_home=tmp_path / "codex-home",
        codex_bin="codex-test",
    )
    calls = []

    def fake_run_codex_agent(invocation, message: str):
        calls.append((invocation, message))
        return subprocess.CompletedProcess(
            args=["codex-test"], returncode=0, stdout="", stderr="runner inspected only"
        )

    def fake_validate_candidate(received_task_dir: Path, attempt: int):
        return _write_validator_report(received_task_dir, attempt=attempt, page_count=1)

    monkeypatch.setattr(workflow, "load_config", lambda: config, raising=False)
    monkeypatch.setattr(workflow, "run_codex_agent", fake_run_codex_agent, raising=False)
    monkeypatch.setattr(
        workflow, "validate_candidate", fake_validate_candidate, raising=False
    )

    run_pdf_to_ppt(task_dir, mode="repair")

    assert [call[0].role for call in calls] == ["runner", "validator"]
    assert "candidate.v2.pptx" in calls[0][1]
    runner_log = (task_dir / "logs" / "runner-repair.log").read_text(encoding="utf-8")
    validator_log = (task_dir / "logs" / "validator-repair.log").read_text(encoding="utf-8")
    assert "runner_missing_artifacts" in runner_log
    assert "deterministic runner fallback" in runner_log
    assert "deterministic validator fallback" in validator_log
    assert (task_dir / "slides" / "slide-model.v2.json").is_file()
    assert (task_dir / "output" / "candidate.v2.pptx").is_file()
    assert (task_dir / "reports" / "runner-repair.v2.json").is_file()
    assert (task_dir / "reports" / "validator.v2.json").is_file()


def test_run_pdf_to_ppt_repair_falls_back_after_validator_success_without_new_report(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
):
    task_dir = tmp_path / "task_1"
    task_dir.mkdir()
    _write_minimal_slide_model(task_dir, attempt=1)
    _write_validator_report(task_dir, attempt=1, page_count=1)
    config = WorkerConfig(
        redis_url="redis://example",
        shared_tasks_dir=tmp_path,
        codex_home=tmp_path / "codex-home",
        codex_bin="codex-test",
    )
    calls = []

    def fake_run_codex_agent(invocation, message: str):
        calls.append((invocation, message))
        if invocation.role == "runner":
            _write_minimal_slide_model(task_dir, attempt=2)
            (task_dir / "output").mkdir(parents=True, exist_ok=True)
            (task_dir / "output" / "candidate.v2.pptx").write_bytes(b"pptx")
            (task_dir / "reports" / "runner-repair.v2.json").write_text(
                "{}",
                encoding="utf-8",
            )
            return subprocess.CompletedProcess(
                args=["codex-test"], returncode=0, stdout="", stderr="runner wrote v2"
            )
        return subprocess.CompletedProcess(
            args=["codex-test"], returncode=0, stdout="", stderr="validator inspected only"
        )

    def fake_validate_candidate(received_task_dir: Path, attempt: int):
        return _write_validator_report(received_task_dir, attempt=attempt, page_count=1)

    monkeypatch.setattr(workflow, "load_config", lambda: config, raising=False)
    monkeypatch.setattr(workflow, "run_codex_agent", fake_run_codex_agent, raising=False)
    monkeypatch.setattr(
        workflow, "validate_candidate", fake_validate_candidate, raising=False
    )

    run_pdf_to_ppt(task_dir, mode="repair")

    assert [call[0].role for call in calls] == ["runner", "validator"]
    assert "validator.v2.json" in calls[1][1]
    assert (task_dir / "logs" / "runner-repair.log").is_file()
    assert (task_dir / "logs" / "validator-repair.log").is_file()
    validator_log = (task_dir / "logs" / "validator-repair.log").read_text(encoding="utf-8")
    assert "validator_missing_report" in validator_log
    assert "deterministic validator fallback" in validator_log
    assert (task_dir / "reports" / "validator.v2.json").is_file()
