from pathlib import Path
from zipfile import ZipFile

import pytest
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

from autofacodex.contracts import SlideModel
from autofacodex.tools.generate_pptx_from_model import generate_from_model
from autofacodex.tools.slide_model_builder import build_initial_slide_model
from autofacodex.tools.pptx_generate import generate_pptx


def _model(slides: list[dict]) -> SlideModel:
    return SlideModel(slides=slides)


def test_generate_pptx_contains_expected_slides_and_editable_text(tmp_path: Path):
    model = _model(
        [
            {
                "page_number": 1,
                "size": {"width": 10, "height": 7.5},
                "elements": [
                    {
                        "id": "title",
                        "type": "text",
                        "text": "Editable Title",
                        "x": 1,
                        "y": 1,
                        "w": 5,
                        "h": 1,
                        "style": {"font_size": 28},
                    }
                ],
                "raster_fallback_regions": [],
            },
            {
                "page_number": 2,
                "size": {"width": 10, "height": 7.5},
                "elements": [],
                "raster_fallback_regions": [],
            },
        ]
    )
    output = tmp_path / "candidate.pptx"

    generate_pptx(model, output)

    presentation = Presentation(output)
    assert len(presentation.slides) == len(model.slides)
    assert presentation.slides[0].shapes[0].text == "Editable Title"
    assert presentation.slides[0].shapes[0].text_frame.paragraphs[0].runs[0].font.size.pt == 28

    with ZipFile(output) as pptx:
        slide_xml = pptx.read("ppt/slides/slide1.xml").decode("utf-8")
    assert "<a:t>Editable Title</a:t>" in slide_xml
    assert "<p:pic>" not in slide_xml


def test_generate_pptx_normalizes_pdf_cjk_font_family(tmp_path: Path):
    model = _model(
        [
            {
                "page_number": 1,
                "size": {"width": 10, "height": 7.5},
                "elements": [
                    {
                        "id": "title",
                        "type": "text",
                        "text": "中文标题",
                        "x": 1,
                        "y": 1,
                        "w": 5,
                        "h": 1,
                        "style": {
                            "font_family": "MicrosoftYaHei-Bold",
                            "font_size": 28,
                            "bold": True,
                        },
                    }
                ],
                "raster_fallback_regions": [],
            }
        ]
    )
    output = tmp_path / "cjk-font.pptx"

    generate_pptx(model, output)

    with ZipFile(output) as archive:
        slide_xml = archive.read("ppt/slides/slide1.xml").decode("utf-8")
    assert 'typeface="Microsoft YaHei"' in slide_xml
    assert "MicrosoftYaHei-Bold" not in slide_xml


def test_generate_pptx_applies_text_rotation(tmp_path: Path):
    model = _model(
        [
            {
                "page_number": 1,
                "size": {"width": 10, "height": 7.5},
                "elements": [
                    {
                        "id": "watermark",
                        "type": "text",
                        "text": "仅供参考",
                        "x": 1,
                        "y": 1,
                        "w": 4,
                        "h": 1,
                        "style": {"font_size": 24, "rotation": -45},
                    }
                ],
                "raster_fallback_regions": [],
            }
        ]
    )
    output = tmp_path / "rotated-text.pptx"

    generate_pptx(model, output)

    shape = Presentation(output).slides[0].shapes[0]
    assert shape.rotation == 315


def test_generate_pptx_preserves_text_style_image_and_rectangle(tmp_path: Path):
    image_path = tmp_path / "assets" / "logo.png"
    image_path.parent.mkdir()
    Image.new("RGB", (12, 10), color=(220, 20, 60)).save(image_path)
    model = _model(
        [
            {
                "page_number": 1,
                "size": {"width": 13.333, "height": 7.5},
                "elements": [
                    {
                        "id": "title",
                        "type": "text",
                        "text": "Styled Title",
                        "x": 1,
                        "y": 0.5,
                        "w": 3,
                        "h": 0.4,
                        "style": {
                            "font_size": 24,
                            "font_family": "Helvetica",
                            "color": "#336699",
                            "bold": True,
                        },
                    },
                    {
                        "id": "image-1",
                        "type": "image",
                        "source": "assets/logo.png",
                        "x": 4,
                        "y": 1,
                        "w": 1.2,
                        "h": 1,
                    },
                    {
                        "id": "rect-1",
                        "type": "shape",
                        "x": 6,
                        "y": 1,
                        "w": 2,
                        "h": 0.8,
                        "style": {
                            "shape": "rect",
                            "line_color": "#112233",
                            "line_width": 1.5,
                            "fill_color": "#FFFFFF",
                        },
                    },
                ],
                "raster_fallback_regions": [],
            }
        ]
    )
    output = tmp_path / "output" / "candidate.pptx"

    generate_pptx(model, output, asset_root=tmp_path)

    slide = Presentation(output).slides[0]
    text_shape = slide.shapes[0]
    run = text_shape.text_frame.paragraphs[0].runs[0]
    assert run.font.size.pt == pytest.approx(24)
    assert run.font.name == "Helvetica"
    assert run.font.bold is True
    assert run.font.color.rgb.__str__() == "336699"
    assert any(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in slide.shapes)
    assert any(
        shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and shape.name.startswith("Rectangle")
        for shape in slide.shapes
    )


def test_generate_pptx_preserves_rich_text_runs(tmp_path: Path):
    model = _model(
        [
            {
                "page_number": 1,
                "size": {"width": 13.333, "height": 7.5},
                "elements": [
                    {
                        "id": "mixed-line",
                        "type": "text",
                        "text": "Hello世界",
                        "x": 1,
                        "y": 1,
                        "w": 3,
                        "h": 0.5,
                        "style": {
                            "runs": [
                                {
                                    "text": "Hello",
                                    "font_size": 16,
                                    "font_family": "Arial",
                                    "color": "#112233",
                                },
                                {
                                    "text": "世界",
                                    "font_size": 18,
                                    "font_family": "KaiTi",
                                    "color": "#445566",
                                },
                            ]
                        },
                    }
                ],
                "raster_fallback_regions": [],
            }
        ]
    )
    output = tmp_path / "rich.pptx"

    generate_pptx(model, output)

    runs = Presentation(output).slides[0].shapes[0].text_frame.paragraphs[0].runs
    assert [run.text for run in runs] == ["Hello", "世界"]
    assert runs[0].font.name == "Arial"
    assert runs[1].font.name == "KaiTi"


def test_generate_pptx_preserves_editable_line_shape(tmp_path: Path):
    model = _model(
        [
            {
                "page_number": 1,
                "size": {"width": 13.333, "height": 7.5},
                "elements": [
                    {
                        "id": "line-1",
                        "type": "shape",
                        "x": 1,
                        "y": 1,
                        "w": 3,
                        "h": 2,
                        "style": {
                            "shape": "line",
                            "x1": 1,
                            "y1": 1,
                            "x2": 4,
                            "y2": 3,
                            "line_color": "#AA0000",
                            "line_width": 2,
                        },
                    }
                ],
                "raster_fallback_regions": [],
            }
        ]
    )
    output = tmp_path / "candidate.pptx"

    generate_pptx(model, output)

    slide = Presentation(output).slides[0]
    assert any(shape.shape_type == MSO_SHAPE_TYPE.LINE for shape in slide.shapes)


def test_generate_pptx_preserves_editable_table(tmp_path: Path):
    model = _model(
        [
            {
                "page_number": 1,
                "size": {"width": 13.333, "height": 7.5},
                "elements": [
                    {
                        "id": "table-1",
                        "type": "table",
                        "x": 1,
                        "y": 1,
                        "w": 5,
                        "h": 1.2,
                        "style": {
                            "font_size": 12,
                            "col_widths": [1.5, 3.5],
                            "row_heights": [0.4, 0.8],
                            "rows": [
                                ["Metric", "Value"],
                                ["ARR", "$12M"],
                            ],
                        },
                    }
                ],
                "raster_fallback_regions": [],
            }
        ]
    )
    output = tmp_path / "table.pptx"

    generate_pptx(model, output)

    slide = Presentation(output).slides[0]
    tables = [shape.table for shape in slide.shapes if shape.has_table]
    assert len(tables) == 1
    table = tables[0]
    assert table.cell(0, 0).text == "Metric"
    assert table.cell(0, 1).text == "Value"
    assert table.cell(1, 0).text == "ARR"
    assert table.cell(1, 1).text == "$12M"
    assert table.columns[0].width == Inches(1.5)
    assert table.columns[1].width == Inches(3.5)
    assert abs(table.rows[0].height - Inches(0.4)) <= 1
    assert abs(table.rows[1].height - Inches(0.8)) <= 1

    with ZipFile(output) as archive:
        slide_xml = archive.read("ppt/slides/slide1.xml").decode("utf-8")
    assert "<a:tbl>" in slide_xml
    assert "<a:t>Metric</a:t>" in slide_xml


def test_generate_pptx_preserves_editable_path(tmp_path: Path):
    model = _model(
        [
            {
                "page_number": 1,
                "size": {"width": 13.333, "height": 7.5},
                "elements": [
                    {
                        "id": "path-1",
                        "type": "path",
                        "x": 1,
                        "y": 1,
                        "w": 2,
                        "h": 1,
                        "style": {
                            "points": [[0, 0], [1, 0], [1, 1], [0, 1]],
                            "closed": True,
                            "fill_color": "#DDEEFF",
                            "line_color": "#336699",
                            "line_width": 1,
                        },
                    }
                ],
                "raster_fallback_regions": [],
            }
        ]
    )
    output = tmp_path / "path.pptx"

    generate_pptx(model, output)

    with ZipFile(output) as archive:
        slide_xml = archive.read("ppt/slides/slide1.xml").decode("utf-8")
    assert "<a:custGeom>" in slide_xml
    assert "336699" in slide_xml
    assert "DDEEFF" in slide_xml


def test_generate_pptx_can_hide_editable_text_overlay(tmp_path: Path):
    model = _model(
        [
            {
                "page_number": 1,
                "size": {"width": 13.333, "height": 7.5},
                "elements": [
                    {
                        "id": "hidden-text",
                        "type": "text",
                        "text": "Editable but visually hidden",
                        "x": 1,
                        "y": 1,
                        "w": 4,
                        "h": 0.4,
                        "style": {
                            "font_size": 18,
                            "color": "#112233",
                            "opacity": 0,
                        },
                    }
                ],
                "raster_fallback_regions": [],
            }
        ]
    )
    output = tmp_path / "candidate.pptx"

    generate_pptx(model, output)

    with ZipFile(output) as archive:
        slide_xml = archive.read("ppt/slides/slide1.xml").decode("utf-8")
    assert "Editable but visually hidden" in slide_xml
    assert 'val="0"' in slide_xml


def test_generate_pptx_applies_editable_shape_opacity(tmp_path: Path):
    model = _model(
        [
            {
                "page_number": 1,
                "size": {"width": 13.333, "height": 7.5},
                "elements": [
                    {
                        "id": "overlay",
                        "type": "shape",
                        "x": 0,
                        "y": 0,
                        "w": 13.333,
                        "h": 7.5,
                        "style": {
                            "shape": "rect",
                            "fill_color": "#000000",
                            "fill_opacity": 0.36863,
                            "line_color": "#000000",
                            "line_opacity": 0.5,
                            "line_width": 1,
                        },
                    }
                ],
                "raster_fallback_regions": [],
            }
        ]
    )
    output = tmp_path / "shape-opacity.pptx"

    generate_pptx(model, output)

    with ZipFile(output) as archive:
        slide_xml = archive.read("ppt/slides/slide1.xml").decode("utf-8")
    assert 'val="36863"' in slide_xml
    assert 'val="50000"' in slide_xml


def test_generate_from_model_uses_task_relative_assets(tmp_path: Path):
    task_dir = tmp_path / "task_1"
    image_path = task_dir / "extracted" / "objects" / "images" / "logo.png"
    model_path = task_dir / "slides" / "slide-model.v2.json"
    output_path = task_dir / "output" / "candidate.v2.pptx"
    image_path.parent.mkdir(parents=True)
    model_path.parent.mkdir(parents=True)
    Image.new("RGB", (12, 10), color=(220, 20, 60)).save(image_path)
    model = _model(
        [
            {
                "page_number": 1,
                "size": {"width": 13.333, "height": 7.5},
                "elements": [
                    {
                        "id": "image-1",
                        "type": "image",
                        "source": "extracted/objects/images/logo.png",
                        "x": 1,
                        "y": 1,
                        "w": 1,
                        "h": 1,
                    }
                ],
                "raster_fallback_regions": [],
            }
        ]
    )
    model_path.write_text(model.model_dump_json(indent=2), encoding="utf-8")

    generate_from_model(model_path, output_path)

    slide = Presentation(output_path).slides[0]
    assert any(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in slide.shapes)


def test_generate_pptx_rejects_mixed_slide_sizes(tmp_path: Path):
    model = _model(
        slides=[
            {
                "page_number": 1,
                "size": {"width": 10, "height": 7.5},
                "elements": [],
                "raster_fallback_regions": [],
            },
            {
                "page_number": 2,
                "size": {"width": 13.333, "height": 7.5},
                "elements": [],
                "raster_fallback_regions": [],
            },
        ]
    )

    with pytest.raises(ValueError, match="same size"):
        generate_pptx(model, tmp_path / "mixed.pptx")


def test_build_initial_slide_model_keeps_empty_slide_for_no_text_page():
    model = build_initial_slide_model(
        {
            "pages": [
                {"page_number": 1, "width": 200, "height": 100, "text": ""},
            ]
        }
    )

    assert len(model.slides) == 1
    assert model.slides[0].elements == []


def test_build_initial_slide_model_normalizes_slide_sizes_across_pages():
    model = build_initial_slide_model(
        {
            "pages": [
                {"page_number": 1, "width": 200, "height": 100, "text": "First"},
                {"page_number": 2, "width": 100, "height": 200, "text": "Second"},
            ]
        }
    )

    assert model.slides[0].size.width == pytest.approx(13.333)
    assert model.slides[0].size.height == pytest.approx(6.6665)
    assert model.slides[1].size == model.slides[0].size


def test_build_initial_slide_model_uses_pdf_blocks_for_positioned_elements():
    model = build_initial_slide_model(
        {
            "pages": [
                {
                    "page_number": 1,
                    "width": 960,
                    "height": 540,
                    "text": "Positioned Title",
                    "text_blocks": [
                        {
                            "type": "text",
                            "bbox": [72, 48, 260, 80],
                            "lines": [
                                {
                                    "bbox": [72, 48, 260, 80],
                                    "spans": [
                                        {
                                            "text": "Positioned Title",
                                            "bbox": [72, 48, 260, 80],
                                            "font": "Helvetica-Bold",
                                            "size": 24,
                                            "color": 0x336699,
                                        }
                                    ],
                                }
                            ],
                        },
                        {
                            "type": "image",
                            "bbox": [288, 108, 384, 216],
                            "source": "extracted/objects/images/page-001-image-001.png",
                        },
                    ],
                    "drawings": [
                        {
                            "shape": "rect",
                            "bbox": [432, 120, 576, 204],
                            "stroke": "#112233",
                            "fill": "#FFFFFF",
                            "stroke_width": 2,
                        }
                    ],
                }
            ]
        }
    )

    slide = model.slides[0]
    assert len(slide.elements) == 3
    text = next(element for element in slide.elements if element.type == "text")
    image = next(element for element in slide.elements if element.type == "image")
    rect = next(element for element in slide.elements if element.type == "shape")
    assert text.type == "text"
    assert text.x == pytest.approx(1.0, abs=0.01)
    assert text.y == pytest.approx(0.667, abs=0.01)
    assert text.w == pytest.approx(2.611, abs=0.01)
    assert text.h == pytest.approx(0.444, abs=0.01)
    assert text.style["font_size"] == pytest.approx(24, abs=0.01)
    assert text.style["font_family"] == "Helvetica-Bold"
    assert text.style["color"] == "#336699"
    assert text.style["bold"] is True
    assert image.type == "image"
    assert image.source == "extracted/objects/images/page-001-image-001.png"
    assert image.x == pytest.approx(4.0, abs=0.01)
    assert rect.type == "shape"
    assert rect.style["shape"] == "rect"
    assert rect.style["line_color"] == "#112233"
    assert rect.style["fill_color"] == "#FFFFFF"


def test_build_initial_slide_model_groups_spans_on_same_line_into_rich_text():
    model = build_initial_slide_model(
        {
            "pages": [
                {
                    "page_number": 1,
                    "width": 960,
                    "height": 540,
                    "text": "Hello世界",
                    "text_blocks": [
                        {
                            "type": "text",
                            "bbox": [72, 48, 260, 80],
                            "lines": [
                                {
                                    "bbox": [72, 48, 260, 80],
                                    "spans": [
                                        {
                                            "text": "Hello",
                                            "bbox": [72, 48, 132, 80],
                                            "font": "ArialMT",
                                            "size": 16,
                                            "color": 0x112233,
                                        },
                                        {
                                            "text": "世界",
                                            "bbox": [132, 48, 260, 80],
                                            "font": "KaiTi",
                                            "size": 18,
                                            "color": 0x445566,
                                        },
                                    ],
                                }
                            ],
                        }
                    ],
                    "drawings": [],
                }
            ]
        }
    )

    [text] = model.slides[0].elements
    assert text.text == "Hello世界"
    assert text.style["runs"][0]["text"] == "Hello"
    assert text.style["runs"][1]["text"] == "世界"
    assert text.style["runs"][0]["font_family"] == "ArialMT"
    assert text.style["runs"][1]["font_family"] == "KaiTi"


def test_build_initial_slide_model_preserves_rotated_text_direction():
    model = build_initial_slide_model(
        {
            "pages": [
                {
                    "page_number": 1,
                    "width": 960,
                    "height": 540,
                    "text": "Tilted Label",
                    "text_blocks": [
                        {
                            "type": "text",
                            "bbox": [100, 100, 400, 400],
                            "lines": [
                                {
                                    "bbox": [100, 100, 400, 400],
                                    "dir": [0.70710678, -0.70710678],
                                    "spans": [
                                        {
                                            "text": "Tilted Label",
                                            "bbox": [100, 100, 400, 400],
                                            "size": 40,
                                        }
                                    ],
                                }
                            ],
                        }
                    ],
                    "drawings": [],
                }
            ]
        }
    )

    [text] = model.slides[0].elements
    assert text.style["rotation"] == pytest.approx(-45)
    assert text.style.get("role") != "watermark"


def test_build_initial_slide_model_hides_confidentiality_watermark_text():
    model = build_initial_slide_model(
        {
            "pages": [
                {
                    "page_number": 1,
                    "width": 960,
                    "height": 540,
                    "text": "仅供隐山资本参考",
                    "text_blocks": [
                        {
                            "type": "text",
                            "bbox": [100, 100, 400, 400],
                            "lines": [
                                {
                                    "bbox": [100, 100, 400, 400],
                                    "dir": [0.70710678, -0.70710678],
                                    "spans": [
                                        {
                                            "text": "仅供隐山资本参考",
                                            "bbox": [100, 100, 400, 400],
                                            "size": 40,
                                            "color": 0xBFBFBF,
                                        }
                                    ],
                                }
                            ],
                        }
                    ],
                    "drawings": [],
                }
            ]
        }
    )

    [text] = model.slides[0].elements
    assert text.style["role"] == "watermark"
    assert text.style["opacity"] == 0


def test_build_initial_slide_model_places_shapes_behind_text_and_images():
    model = build_initial_slide_model(
        {
            "pages": [
                {
                    "page_number": 1,
                    "width": 960,
                    "height": 540,
                    "text": "Title",
                    "text_blocks": [
                        {
                            "type": "text",
                            "bbox": [72, 48, 260, 80],
                            "lines": [
                                {
                                    "bbox": [72, 48, 260, 80],
                                    "spans": [
                                        {
                                            "text": "Title",
                                            "bbox": [72, 48, 260, 80],
                                            "font": "Helvetica",
                                            "size": 24,
                                            "color": 0,
                                        }
                                    ],
                                }
                            ],
                        },
                        {
                            "type": "image",
                            "bbox": [288, 108, 384, 216],
                            "source": "extracted/objects/images/page-001-image-001.png",
                        },
                    ],
                    "drawings": [
                        {
                            "shape": "rect",
                            "bbox": [0, 0, 960, 540],
                            "fill": "#FFFFFF",
                        }
                    ],
                }
            ]
        }
    )

    assert [element.type for element in model.slides[0].elements] == [
        "shape",
        "text",
        "image",
    ]
    background = model.slides[0].elements[0]
    assert background.style["line_width"] == 0
    assert background.style["line_color"] == "#FFFFFF"
    assert background.style["role"] == "background"


def test_build_initial_slide_model_preserves_shape_opacity():
    model = build_initial_slide_model(
        {
            "pages": [
                {
                    "page_number": 1,
                    "width": 960,
                    "height": 540,
                    "text": "",
                    "text_blocks": [],
                    "drawings": [
                        {
                            "shape": "rect",
                            "bbox": [0, 0, 960, 540],
                            "fill": "#000000",
                            "stroke": "#000000",
                            "stroke_width": 1,
                            "fill_opacity": 0.36863,
                            "stroke_opacity": 0.5,
                            "seqno": 8,
                        }
                    ],
                }
            ]
        }
    )

    [overlay] = model.slides[0].elements
    assert overlay.style["fill_opacity"] == pytest.approx(0.36863)
    assert overlay.style["line_opacity"] == pytest.approx(0.5)


def test_build_initial_slide_model_marks_full_page_images_as_background():
    model = build_initial_slide_model(
        {
            "pages": [
                {
                    "page_number": 1,
                    "width": 960,
                    "height": 540,
                    "text": "Title",
                    "text_blocks": [
                        {
                            "type": "image",
                            "bbox": [0, 0, 960, 540],
                            "source": "objects/images/page-001-image-001.png",
                            "seqno": 1,
                        },
                        {
                            "type": "text",
                            "bbox": [72, 48, 260, 80],
                            "lines": [
                                {
                                    "bbox": [72, 48, 260, 80],
                                    "spans": [
                                        {
                                            "text": "Title",
                                            "bbox": [72, 48, 260, 80],
                                            "font": "Helvetica",
                                            "size": 24,
                                            "color": 0,
                                            "seqno": 2,
                                        }
                                    ],
                                }
                            ],
                        },
                    ],
                    "drawings": [],
                }
            ]
        }
    )

    background = model.slides[0].elements[0]
    assert background.type == "image"
    assert background.source == "extracted/objects/images/page-001-image-001.png"
    assert background.style["role"] == "background"


def test_build_initial_slide_model_collapses_grid_region_into_editable_table():
    model = build_initial_slide_model(
        {
            "pages": [
                {
                    "page_number": 1,
                    "width": 960,
                    "height": 540,
                    "text": "排序名称数量1千帆星座150002GW星座12992",
                    "text_blocks": [
                        {
                            "type": "text",
                            "bbox": [80, 220, 420, 300],
                            "lines": [
                                {
                                    "bbox": [88, 225, 112, 241],
                                    "spans": [{"text": "排序", "bbox": [88, 225, 112, 241], "size": 12}],
                                },
                                {
                                    "bbox": [170, 225, 205, 241],
                                    "spans": [{"text": "名称", "bbox": [170, 225, 205, 241], "size": 12}],
                                },
                                {
                                    "bbox": [374, 225, 410, 241],
                                    "spans": [{"text": "数量", "bbox": [374, 225, 410, 241], "size": 12}],
                                },
                                {
                                    "bbox": [97, 252, 104, 267],
                                    "spans": [{"text": "1", "bbox": [97, 252, 104, 267], "size": 12}],
                                },
                                {
                                    "bbox": [168, 252, 217, 267],
                                    "spans": [{"text": "千帆星座", "bbox": [168, 252, 217, 267], "size": 12}],
                                },
                                {
                                    "bbox": [375, 252, 410, 267],
                                    "spans": [{"text": "15000", "bbox": [375, 252, 410, 267], "size": 12}],
                                },
                                {
                                    "bbox": [97, 277, 104, 293],
                                    "spans": [{"text": "2", "bbox": [97, 277, 104, 293], "size": 12}],
                                },
                                {
                                    "bbox": [170, 277, 215, 293],
                                    "spans": [{"text": "GW星座", "bbox": [170, 277, 215, 293], "size": 12}],
                                },
                                {
                                    "bbox": [375, 277, 410, 293],
                                    "spans": [{"text": "12992", "bbox": [375, 277, 410, 293], "size": 12}],
                                },
                                {
                                    "bbox": [60, 230, 360, 320],
                                    "spans": [
                                        {
                                            "text": "内部参考水印",
                                            "bbox": [60, 230, 360, 320],
                                            "size": 12,
                                        }
                                    ],
                                },
                            ],
                        }
                    ],
                    "drawings": [
                        {"shape": "rect", "bbox": [77, 220, 124, 247], "fill": "#DAE3F5", "seqno": 1},
                        {"shape": "rect", "bbox": [124, 220, 262, 247], "fill": "#DAE3F5", "seqno": 2},
                        {"shape": "rect", "bbox": [262, 220, 428, 247], "fill": "#DAE3F5", "seqno": 3},
                        {
                            "shape": "line",
                            "bbox": [124, 220, 124, 300],
                            "p1": [124, 220],
                            "p2": [124, 300],
                            "stroke": "#FFFFFF",
                            "stroke_width": 1,
                            "seqno": 4,
                        },
                        {
                            "shape": "line",
                            "bbox": [262, 220, 262, 300],
                            "p1": [262, 220],
                            "p2": [262, 300],
                            "stroke": "#FFFFFF",
                            "stroke_width": 1,
                            "seqno": 5,
                        },
                        {
                            "shape": "line",
                            "bbox": [77, 247, 428, 247],
                            "p1": [77, 247],
                            "p2": [428, 247],
                            "stroke": "#2E54A1",
                            "stroke_width": 1,
                            "seqno": 6,
                        },
                        {
                            "shape": "line",
                            "bbox": [77, 273, 428, 273],
                            "p1": [77, 273],
                            "p2": [428, 273],
                            "stroke": "#BFBFBF",
                            "stroke_width": 0.5,
                            "seqno": 7,
                        },
                        {
                            "shape": "line",
                            "bbox": [77, 298, 428, 298],
                            "p1": [77, 298],
                            "p2": [428, 298],
                            "stroke": "#BFBFBF",
                            "stroke_width": 0.5,
                            "seqno": 8,
                        },
                    ],
                }
            ]
        }
    )

    table = next(element for element in model.slides[0].elements if element.type == "table")
    assert table.style["rows"] == [
        ["排序", "名称", "数量"],
        ["1", "千帆星座", "15000"],
        ["2", "GW星座", "12992"],
    ]
    assert "内部参考水印" not in "\n".join(cell for row in table.style["rows"] for cell in row)
    assert table.style["covered_text_ids"] == [
        "p1-text-1",
        "p1-text-2",
        "p1-text-3",
        "p1-text-4",
        "p1-text-5",
        "p1-text-6",
        "p1-text-7",
        "p1-text-8",
        "p1-text-9",
    ]
    assert sum(table.style["col_widths"]) == pytest.approx(table.w)
    assert sum(table.style["row_heights"]) == pytest.approx(table.h)
    assert table.style["font_size"] == pytest.approx(12 * 13.333 * 72 / 960)
    assert table.style["role"] == "semantic_table"
    assert table.style["opacity"] == 0
    assert any(element.type == "text" and element.text == "千帆星座" for element in model.slides[0].elements)


def test_build_initial_slide_model_orders_elements_by_pdf_paint_order():
    model = build_initial_slide_model(
        {
            "pages": [
                {
                    "page_number": 1,
                    "width": 960,
                    "height": 540,
                    "text": "Title",
                    "text_blocks": [
                        {
                            "type": "text",
                            "bbox": [72, 48, 260, 80],
                            "lines": [
                                {
                                    "bbox": [72, 48, 260, 80],
                                    "spans": [
                                        {
                                            "text": "Title",
                                            "bbox": [72, 48, 260, 80],
                                            "font": "Helvetica",
                                            "size": 24,
                                            "color": 0,
                                            "seqno": 3,
                                        }
                                    ],
                                }
                            ],
                        },
                        {
                            "type": "image",
                            "bbox": [288, 108, 384, 216],
                            "source": "extracted/objects/images/page-001-image-001.png",
                            "seqno": 2,
                        },
                    ],
                    "drawings": [
                        {
                            "shape": "rect",
                            "bbox": [0, 0, 960, 540],
                            "fill": "#FFFFFF",
                            "seqno": 1,
                        },
                        {
                            "shape": "rect",
                            "bbox": [432, 120, 576, 204],
                            "fill": "#DDDDDD",
                            "seqno": 4,
                        },
                    ],
                }
            ]
        }
    )

    assert [element.id for element in model.slides[0].elements] == [
        "p1-shape-1",
        "p1-image-1",
        "p1-text-1",
        "p1-shape-2",
    ]


def test_build_initial_slide_model_converts_line_drawings_to_editable_shape():
    model = build_initial_slide_model(
        {
            "pages": [
                {
                    "page_number": 1,
                    "width": 400,
                    "height": 300,
                    "text": "",
                    "text_blocks": [],
                    "drawings": [
                        {
                            "shape": "line",
                            "p1": [10, 280],
                            "p2": [110, 180],
                            "bbox": [10, 180, 110, 280],
                            "stroke": "#FF0000",
                            "stroke_width": 2,
                            "seqno": 1,
                        }
                    ],
                }
            ]
        }
    )

    line = model.slides[0].elements[0]
    assert line.type == "shape"
    assert line.style["shape"] == "line"
    assert line.style["x1"] == pytest.approx(0.333, abs=0.01)
    assert line.style["y1"] == pytest.approx(9.333, abs=0.01)
    assert line.style["x2"] == pytest.approx(3.667, abs=0.01)
    assert line.style["y2"] == pytest.approx(6.0, abs=0.01)
    assert line.style["line_color"] == "#FF0000"


def test_build_initial_slide_model_normalizes_extracted_image_sources_to_task_relative():
    model = build_initial_slide_model(
        {
            "pages": [
                {
                    "page_number": 1,
                    "width": 960,
                    "height": 540,
                    "text": "",
                    "text_blocks": [
                        {
                            "type": "image",
                            "bbox": [0, 0, 96, 96],
                            "source": "objects/images/page-001-image-001.png",
                        }
                    ],
                    "drawings": [],
                }
            ]
        }
    )

    assert model.slides[0].elements[0].source == (
        "extracted/objects/images/page-001-image-001.png"
    )


@pytest.mark.parametrize(
    ("page", "message"),
    [
        ({"page_number": 1, "height": 100, "text": "Missing width"}, "width"),
        ({"page_number": 1, "width": 100, "text": "Missing height"}, "height"),
        ({"page_number": 1, "width": 0, "height": 100, "text": "Zero width"}, "width"),
        ({"page_number": 1, "width": 100, "height": 0, "text": "Zero height"}, "height"),
        ({"page_number": 1, "width": -1, "height": 100, "text": "Negative width"}, "width"),
        ({"page_number": 1, "width": 100, "height": -1, "text": "Negative height"}, "height"),
        ({"page_number": 1, "width": "wide", "height": 100, "text": "Bad width"}, "width"),
        ({"page_number": 1, "width": 100, "height": "tall", "text": "Bad height"}, "height"),
    ],
)
def test_build_initial_slide_model_rejects_invalid_dimensions(page: dict, message: str):
    with pytest.raises(ValueError, match=message):
        build_initial_slide_model({"pages": [page]})


def test_build_initial_slide_model_suppresses_fragments_inside_dominant_background_image():
    model = build_initial_slide_model(
        {
            "pages": [
                {
                    "page_number": 1,
                    "width": 960,
                    "height": 540,
                    "text": "Foreground Title",
                    "text_blocks": [
                        {
                            "type": "image",
                            "bbox": [0, 60, 960, 540],
                            "source": "objects/images/page-001-image-001.png",
                            "seqno": 1,
                        },
                        {
                            "type": "image",
                            "bbox": [320, 220, 460, 340],
                            "source": "objects/images/page-001-image-002.png",
                            "seqno": 2,
                        },
                        {
                            "type": "text",
                            "bbox": [72, 80, 360, 110],
                            "lines": [
                                {
                                    "bbox": [72, 80, 360, 110],
                                    "spans": [
                                        {
                                            "text": "Foreground Title",
                                            "bbox": [72, 80, 360, 110],
                                            "font": "Helvetica",
                                            "size": 24,
                                            "color": 0,
                                            "seqno": 4,
                                        }
                                    ],
                                }
                            ],
                        },
                        {
                            "type": "image",
                            "bbox": [20, 20, 80, 80],
                            "source": "objects/images/page-001-image-003.png",
                            "seqno": 5,
                        },
                    ],
                    "drawings": [
                        {
                            "shape": "rect",
                            "bbox": [350, 260, 450, 300],
                            "fill": "#DDDDDD",
                            "stroke": "#DDDDDD",
                            "seqno": 3,
                        }
                    ],
                }
            ]
        }
    )

    elements = model.slides[0].elements
    assert [element.id for element in elements] == [
        "p1-image-1",
        "p1-text-1",
        "p1-image-3",
    ]
    assert elements[0].style["role"] == "background"
    assert elements[1].type == "text"
    assert elements[1].text == "Foreground Title"
    assert elements[2].source == "extracted/objects/images/page-001-image-003.png"


def test_build_initial_slide_model_keeps_fragments_without_dominant_background_image():
    model = build_initial_slide_model(
        {
            "pages": [
                {
                    "page_number": 1,
                    "width": 960,
                    "height": 540,
                    "text": "",
                    "text_blocks": [
                        {
                            "type": "image",
                            "bbox": [0, 60, 400, 300],
                            "source": "objects/images/page-001-image-001.png",
                            "seqno": 1,
                        },
                        {
                            "type": "image",
                            "bbox": [320, 220, 460, 340],
                            "source": "objects/images/page-001-image-002.png",
                            "seqno": 2,
                        },
                    ],
                    "drawings": [],
                }
            ]
        }
    )

    assert [element.id for element in model.slides[0].elements] == [
        "p1-image-1",
        "p1-image-2",
    ]
    assert "role" not in model.slides[0].elements[0].style
