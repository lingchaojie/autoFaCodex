from pathlib import Path

from pptx import Presentation

from autofacodex.evaluation.compare_ideal_pptx import compare_pptx_structure


def _pptx(path: Path, texts: list[str]) -> None:
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for text in texts:
        slide = prs.slides.add_slide(blank)
        box = slide.shapes.add_textbox(914400, 914400, 3657600, 914400)
        box.text = text
    prs.save(path)


def test_compare_pptx_structure_reports_page_count_and_text_delta(tmp_path: Path):
    generated = tmp_path / "generated.pptx"
    ideal = tmp_path / "ideal.pptx"
    _pptx(generated, ["One"])
    _pptx(ideal, ["One", "Two"])

    result = compare_pptx_structure(generated, ideal)

    assert result["generated_slide_count"] == 1
    assert result["ideal_slide_count"] == 2
    assert result["slide_count_delta"] == -1
    assert result["pages"][0]["generated_text_runs"] == 1
    assert result["pages"][0]["ideal_text_runs"] == 1


def test_compare_pptx_structure_includes_strategy_deltas(tmp_path: Path, monkeypatch):
    generated = tmp_path / "generated.pptx"
    ideal = tmp_path / "ideal.pptx"
    _pptx(generated, ["Generated"])
    _pptx(ideal, ["Ideal"])

    profiles = {
        generated: {
            "pages": [
                {
                    "page_number": 1,
                    "strategy": "fragmented_objects",
                    "text_box_count": 8,
                    "pictures": 12,
                    "shapes": 36,
                    "largest_picture_area_ratio": 0.72,
                    "picture_coverage_ratio": 0.89,
                }
            ]
        },
        ideal: {
            "pages": [
                {
                    "page_number": 1,
                    "strategy": "background_plus_foreground_text",
                    "text_box_count": 3,
                    "pictures": 3,
                    "shapes": 8,
                    "largest_picture_area_ratio": 0.94,
                    "picture_coverage_ratio": 0.95,
                }
            ]
        },
    }
    monkeypatch.setattr(
        "autofacodex.evaluation.compare_ideal_pptx.profile_pptx_strategy",
        lambda path: profiles[path],
    )

    result = compare_pptx_structure(generated, ideal)

    page = result["pages"][0]
    assert page["generated_strategy"] == "fragmented_objects"
    assert page["ideal_strategy"] == "background_plus_foreground_text"
    assert page["strategy_matches"] is False
    assert page["picture_count_delta"] == 9
    assert page["shape_count_delta"] == 28
    assert page["text_box_count_delta"] == 5
    assert page["largest_picture_area_ratio_delta"] == -0.22
    assert page["picture_coverage_ratio_delta"] == -0.06


def test_compare_pptx_structure_reports_top_geometry_mismatches(tmp_path: Path, monkeypatch):
    generated = tmp_path / "generated.pptx"
    ideal = tmp_path / "ideal.pptx"
    _pptx(generated, ["Generated"])
    _pptx(ideal, ["Ideal"])

    profiles = {
        generated: {
            "pages": [
                {
                    "page_number": 1,
                    "size": {"width": 10, "height": 5},
                    "picture_geometries": [{"x": 1, "y": 1, "w": 4, "h": 2}],
                    "shape_geometries": [{"x": 0, "y": 0, "w": 1, "h": 1}],
                    "text_box_geometries": [{"x": 2, "y": 2, "w": 2, "h": 1}],
                }
            ]
        },
        ideal: {
            "pages": [
                {
                    "page_number": 1,
                    "size": {"width": 10, "height": 5},
                    "picture_geometries": [{"x": 1, "y": 1, "w": 4, "h": 2}],
                    "shape_geometries": [{"x": 5, "y": 0, "w": 1, "h": 1}],
                    "text_box_geometries": [{"x": 2, "y": 3, "w": 2, "h": 1}],
                }
            ]
        },
    }
    monkeypatch.setattr(
        "autofacodex.evaluation.compare_ideal_pptx.profile_pptx_strategy",
        lambda path: profiles[path],
    )

    result = compare_pptx_structure(generated, ideal)

    mismatches = result["pages"][0]["top_geometry_mismatches"]
    assert mismatches[0] == {
        "kind": "shape",
        "index": 0,
        "score": 0.5,
        "generated_bbox": [0.0, 0.0, 0.1, 0.2],
        "ideal_bbox": [0.5, 0.0, 0.6, 0.2],
    }
    assert mismatches[1]["kind"] == "text_box"
    assert mismatches[1]["score"] == 0.2
