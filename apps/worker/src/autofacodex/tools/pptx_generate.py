from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

from autofacodex.contracts import SlideElement, SlideModel


def _validate_single_slide_size(model: SlideModel) -> None:
    if not model.slides:
        return

    expected_size = model.slides[0].size
    for slide_spec in model.slides[1:]:
        if slide_spec.size != expected_size:
            raise ValueError("All slides must use the same size before PPTX generation")


def _rgb(color: object) -> RGBColor | None:
    if not isinstance(color, str):
        return None
    value = color.strip()
    if value.startswith("#"):
        value = value[1:]
    if len(value) != 6:
        return None
    try:
        int(value, 16)
    except ValueError:
        return None
    return RGBColor.from_string(value.upper())


def _font_family(value: object) -> str | None:
    if not isinstance(value, str):
        return None
    font = value.strip()
    if not font:
        return None
    normalized = (
        font.lower()
        .replace("-", "")
        .replace("_", "")
        .replace(" ", "")
        .replace(",", "")
    )
    mappings = {
        "microsoftyahei": "Microsoft YaHei",
        "microsoftyaheibold": "Microsoft YaHei",
        "microsoftyaheilight": "Microsoft YaHei",
        "microsoftyaheiui": "Microsoft YaHei UI",
        "microsoftyaheiuibold": "Microsoft YaHei UI",
        "microsoftyaheiuilight": "Microsoft YaHei UI",
        "arialmt": "Arial",
        "cambriamath": "Cambria Math",
        "simsun": "SimSun",
        "simhei": "SimHei",
        "dengxian": "DengXian",
    }
    return mappings.get(normalized, font)


def _set_font_typeface(run, family: str) -> None:
    run.font.name = family
    rpr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        node = rpr.find(qn(tag))
        if node is None:
            node = OxmlElement(tag)
            rpr.append(node)
        node.set("typeface", family)


def _opacity(value: object) -> float | None:
    if value is None:
        return None
    try:
        opacity = float(value)
    except (TypeError, ValueError):
        return None
    return max(0.0, min(1.0, opacity))


def _apply_font_opacity(run, opacity: float) -> None:
    if opacity >= 1:
        return
    if run.font.color.rgb is None:
        run.font.color.rgb = RGBColor(0, 0, 0)
    rpr = run._r.get_or_add_rPr()
    solid_fill = rpr.find(qn("a:solidFill"))
    if solid_fill is None:
        return
    srgb = solid_fill.find(qn("a:srgbClr"))
    if srgb is None:
        return
    for existing in srgb.findall(qn("a:alpha")):
        srgb.remove(existing)
    alpha = OxmlElement("a:alpha")
    alpha.set("val", str(round(opacity * 100000)))
    srgb.append(alpha)


def _apply_srgb_opacity(srgb, opacity: float | None) -> None:
    if opacity is None or opacity >= 1:
        return
    for existing in srgb.findall(qn("a:alpha")):
        srgb.remove(existing)
    alpha = OxmlElement("a:alpha")
    alpha.set("val", str(round(opacity * 100000)))
    srgb.append(alpha)


def _apply_shape_fill_opacity(shape, opacity: float | None) -> None:
    if opacity is None or opacity >= 1:
        return
    solid_fill = shape._element.spPr.find(qn("a:solidFill"))
    if solid_fill is None:
        return
    srgb = solid_fill.find(qn("a:srgbClr"))
    if srgb is None:
        return
    _apply_srgb_opacity(srgb, opacity)


def _apply_shape_line_opacity(shape, opacity: float | None) -> None:
    if opacity is None or opacity >= 1:
        return
    line = shape._element.spPr.find(qn("a:ln"))
    if line is None:
        return
    solid_fill = line.find(qn("a:solidFill"))
    if solid_fill is None:
        return
    srgb = solid_fill.find(qn("a:srgbClr"))
    if srgb is None:
        return
    _apply_srgb_opacity(srgb, opacity)


def _resolve_source(source: str | None, output_path: Path, asset_root: Path | None) -> Path:
    if not source:
        raise ValueError("Image element is missing source")
    source_path = Path(source)
    if source_path.is_absolute():
        return source_path
    root = asset_root if asset_root is not None else output_path.parent.parent
    return root / source_path


def _apply_run_style(run, style: dict) -> None:
    font_size = style.get("font_size", 18)
    run.font.size = Pt(float(font_size))
    if family := _font_family(style.get("font_family")):
        _set_font_typeface(run, family)
    if style.get("bold") is not None:
        run.font.bold = bool(style["bold"])
    if style.get("italic") is not None:
        run.font.italic = bool(style["italic"])
    if color := _rgb(style.get("color")):
        run.font.color.rgb = color
    if (opacity := _opacity(style.get("opacity"))) is not None:
        _apply_font_opacity(run, opacity)


def _add_text(slide, element: SlideElement) -> None:
    box = slide.shapes.add_textbox(
        Inches(element.x),
        Inches(element.y),
        Inches(element.w),
        Inches(element.h),
    )
    if element.style.get("rotation") is not None:
        box.rotation = float(element.style["rotation"])
    text_frame = box.text_frame
    text_frame.margin_left = 0
    text_frame.margin_right = 0
    text_frame.margin_top = 0
    text_frame.margin_bottom = 0
    paragraph = text_frame.paragraphs[0]
    style = element.style
    runs = style.get("runs")
    if isinstance(runs, list) and runs:
        for run_style in runs:
            if not isinstance(run_style, dict):
                continue
            run = paragraph.add_run()
            run.text = str(run_style.get("text", ""))
            _apply_run_style(run, {**style, **run_style})
        return
    run = paragraph.add_run()
    run.text = element.text or ""
    _apply_run_style(run, style)


def _add_image(slide, element: SlideElement, output_path: Path, asset_root: Path | None) -> None:
    image_path = _resolve_source(element.source, output_path, asset_root)
    if not image_path.is_file():
        raise FileNotFoundError(f"Image asset not found: {image_path}")
    slide.shapes.add_picture(
        str(image_path),
        Inches(element.x),
        Inches(element.y),
        width=Inches(element.w),
        height=Inches(element.h),
    )


def _add_shape(slide, element: SlideElement) -> None:
    shape_type = element.style.get("shape")
    if shape_type == "line":
        _add_line(slide, element)
        return
    if shape_type != "rect":
        return
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(element.x),
        Inches(element.y),
        Inches(element.w),
        Inches(element.h),
    )
    if color := _rgb(element.style.get("fill_color")):
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        _apply_shape_fill_opacity(shape, _opacity(element.style.get("fill_opacity")))
    else:
        shape.fill.background()
    if color := _rgb(element.style.get("line_color")):
        shape.line.color.rgb = color
        _apply_shape_line_opacity(shape, _opacity(element.style.get("line_opacity")))
    if element.style.get("line_width") is not None:
        shape.line.width = Pt(float(element.style["line_width"]))


def _coordinate(style: dict, name: str, fallback: float) -> float:
    try:
        return float(style.get(name, fallback))
    except (TypeError, ValueError):
        return fallback


def _add_line(slide, element: SlideElement) -> None:
    style = element.style
    x1 = _coordinate(style, "x1", element.x)
    y1 = _coordinate(style, "y1", element.y)
    x2 = _coordinate(style, "x2", element.x + element.w)
    y2 = _coordinate(style, "y2", element.y + element.h)
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    if color := _rgb(style.get("line_color")):
        line.line.color.rgb = color
        _apply_shape_line_opacity(line, _opacity(style.get("line_opacity")))
    if style.get("line_width") is not None:
        line.line.width = Pt(float(style["line_width"]))


def _table_rows(element: SlideElement) -> list[list[object]]:
    rows = element.style.get("rows")
    if not isinstance(rows, list):
        return []
    return [row for row in rows if isinstance(row, list)]


def _cell_text(cell_value: object) -> str:
    if isinstance(cell_value, dict):
        return str(cell_value.get("text", ""))
    return str(cell_value)


def _cell_style(table_style: dict, cell_value: object) -> dict:
    if isinstance(cell_value, dict):
        return {**table_style, **cell_value}
    return table_style


def _table_dimensions(style: dict, key: str, count: int, total: float) -> list[float]:
    if count <= 0:
        return []
    values = style.get(key)
    if isinstance(values, list) and len(values) == count:
        try:
            dimensions = [float(value) for value in values]
        except (TypeError, ValueError):
            dimensions = []
        if dimensions and all(value > 0 for value in dimensions):
            dimension_total = sum(dimensions)
            if dimension_total > 0 and total > 0:
                scale = total / dimension_total
                return [value * scale for value in dimensions]
    return [total / count for _ in range(count)]


def _hide_table_cell_borders(cell) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    for tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        line = tc_pr.find(qn(tag))
        if line is None:
            line = OxmlElement(tag)
            tc_pr.append(line)
        for child in list(line):
            line.remove(child)
        line.append(OxmlElement("a:noFill"))


def _add_table(slide, element: SlideElement) -> None:
    rows = _table_rows(element)
    if not rows:
        return
    row_count = len(rows)
    col_count = max((len(row) for row in rows), default=0)
    if col_count == 0:
        return
    graphic_frame = slide.shapes.add_table(
        row_count,
        col_count,
        Inches(element.x),
        Inches(element.y),
        Inches(element.w),
        Inches(element.h),
    )
    table = graphic_frame.table
    for col_index, width in enumerate(
        _table_dimensions(element.style, "col_widths", col_count, element.w)
    ):
        table.columns[col_index].width = Inches(width)
    for row_index, height in enumerate(
        _table_dimensions(element.style, "row_heights", row_count, element.h)
    ):
        table.rows[row_index].height = Inches(height)
    for row_index, row in enumerate(rows):
        for col_index in range(col_count):
            cell_value = row[col_index] if col_index < len(row) else ""
            cell = table.cell(row_index, col_index)
            cell.fill.background()
            _hide_table_cell_borders(cell)
            cell.margin_left = 0
            cell.margin_right = 0
            cell.margin_top = 0
            cell.margin_bottom = 0
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.text = _cell_text(cell_value)
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.CENTER
            if paragraph.runs:
                _apply_run_style(
                    paragraph.runs[0],
                    _cell_style(element.style, cell_value),
                )


def _path_points(element: SlideElement) -> list[tuple[float, float]]:
    points = element.style.get("points")
    if not isinstance(points, list):
        return []
    normalized_points: list[tuple[float, float]] = []
    for point in points:
        if not isinstance(point, (list, tuple)) or len(point) != 2:
            return []
        try:
            x = float(point[0])
            y = float(point[1])
        except (TypeError, ValueError):
            return []
        normalized_points.append((x * 1000, y * 1000))
    return normalized_points


def _add_path(slide, element: SlideElement) -> None:
    points = _path_points(element)
    if len(points) < 2 or element.w <= 0 or element.h <= 0:
        return
    builder = slide.shapes.build_freeform(
        points[0][0],
        points[0][1],
        scale=(Inches(element.w) / 1000, Inches(element.h) / 1000),
    )
    builder.add_line_segments(points[1:], close=bool(element.style.get("closed", True)))
    shape = builder.convert_to_shape(Inches(element.x), Inches(element.y))
    if color := _rgb(element.style.get("fill_color")):
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        _apply_shape_fill_opacity(shape, _opacity(element.style.get("fill_opacity")))
    else:
        shape.fill.background()
    if color := _rgb(element.style.get("line_color")):
        shape.line.color.rgb = color
        _apply_shape_line_opacity(shape, _opacity(element.style.get("line_opacity")))
    if element.style.get("line_width") is not None:
        shape.line.width = Pt(float(element.style["line_width"]))


def generate_pptx(
    model: SlideModel, output_path: Path, asset_root: Path | None = None
) -> Path:
    output_path.parent.mkdir(parents=True, exist_ok=True)
    _validate_single_slide_size(model)
    presentation = Presentation()
    if model.slides:
        presentation.slide_width = Inches(model.slides[0].size.width)
        presentation.slide_height = Inches(model.slides[0].size.height)

    blank_layout = presentation.slide_layouts[6]
    for slide_spec in model.slides:
        slide = presentation.slides.add_slide(blank_layout)
        for element in slide_spec.elements:
            if element.type == "text":
                _add_text(slide, element)
            elif element.type == "image":
                _add_image(slide, element, output_path, asset_root)
            elif element.type == "shape":
                _add_shape(slide, element)
            elif element.type == "table":
                _add_table(slide, element)
            elif element.type == "path":
                _add_path(slide, element)
    presentation.save(output_path)
    return output_path
