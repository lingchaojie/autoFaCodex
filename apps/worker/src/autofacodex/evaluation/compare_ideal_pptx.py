import math
from pathlib import Path

from pptx import Presentation

from autofacodex.evaluation.pptx_strategy import profile_pptx_strategy


def _slide_counts(slide) -> dict[str, int]:
    text_runs = 0
    pictures = 0
    shapes = 0
    tables = 0
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
            text_runs += 1
        if shape.shape_type == 13:
            pictures += 1
        if getattr(shape, "has_table", False):
            tables += 1
        shapes += 1
    return {
        "text_runs": text_runs,
        "pictures": pictures,
        "shapes": shapes,
        "tables": tables,
    }


def _profile_by_page(profile: dict) -> dict[int, dict]:
    return {
        int(page["page_number"]): page
        for page in profile.get("pages", [])
        if page.get("page_number") is not None
    }


def _profile_number(page: dict, field: str) -> float:
    try:
        number = float(page.get(field, 0) or 0)
    except (OverflowError, TypeError, ValueError):
        return 0.0
    if not math.isfinite(number):
        return 0.0
    return number


def _profile_int(page: dict, field: str) -> int:
    return int(_profile_number(page, field))


def _geometry_list(page: dict, field: str) -> list[dict]:
    geometries = page.get(field, [])
    if not isinstance(geometries, list):
        return []
    return [geometry for geometry in geometries if isinstance(geometry, dict)]


def _normalized_bbox(geometry: dict | None, size: dict) -> list[float] | None:
    if geometry is None:
        return None
    width = _profile_number(size, "width")
    height = _profile_number(size, "height")
    if width <= 0 or height <= 0:
        return None
    try:
        x = float(geometry.get("x", 0) or 0)
        y = float(geometry.get("y", 0) or 0)
        w = float(geometry.get("w", 0) or 0)
        h = float(geometry.get("h", 0) or 0)
    except (OverflowError, TypeError, ValueError):
        return None
    if not all(math.isfinite(value) for value in (x, y, w, h)):
        return None
    return [
        round(x / width, 6),
        round(y / height, 6),
        round((x + w) / width, 6),
        round((y + h) / height, 6),
    ]


def _bbox_mismatch_score(
    generated_bbox: list[float] | None, ideal_bbox: list[float] | None
) -> float:
    if generated_bbox is None or ideal_bbox is None:
        return 1.0
    return round(max(abs(left - right) for left, right in zip(generated_bbox, ideal_bbox)), 6)


def _top_geometry_mismatches(generated_page: dict, ideal_page: dict) -> list[dict]:
    fields = (
        ("picture", "picture_geometries"),
        ("shape", "shape_geometries"),
        ("text_box", "text_box_geometries"),
    )
    generated_size = (
        generated_page.get("size") if isinstance(generated_page.get("size"), dict) else {}
    )
    ideal_size = ideal_page.get("size") if isinstance(ideal_page.get("size"), dict) else {}
    mismatches = []
    for kind, field in fields:
        generated_geometries = _geometry_list(generated_page, field)
        ideal_geometries = _geometry_list(ideal_page, field)
        for index in range(max(len(generated_geometries), len(ideal_geometries))):
            generated_bbox = _normalized_bbox(
                generated_geometries[index] if index < len(generated_geometries) else None,
                generated_size,
            )
            ideal_bbox = _normalized_bbox(
                ideal_geometries[index] if index < len(ideal_geometries) else None,
                ideal_size,
            )
            if generated_bbox == ideal_bbox:
                continue
            mismatches.append(
                {
                    "kind": kind,
                    "index": index,
                    "score": _bbox_mismatch_score(generated_bbox, ideal_bbox),
                    "generated_bbox": generated_bbox,
                    "ideal_bbox": ideal_bbox,
                }
            )
    return sorted(
        mismatches,
        key=lambda item: (-item["score"], item["kind"], item["index"]),
    )[:5]


def _strategy_delta(generated_page: dict, ideal_page: dict) -> dict:
    return {
        "generated_strategy": generated_page.get("strategy", "unknown"),
        "ideal_strategy": ideal_page.get("strategy", "unknown"),
        "strategy_matches": generated_page.get("strategy") == ideal_page.get("strategy"),
        "picture_count_delta": _profile_int(generated_page, "pictures")
        - _profile_int(ideal_page, "pictures"),
        "shape_count_delta": _profile_int(generated_page, "shapes")
        - _profile_int(ideal_page, "shapes"),
        "text_box_count_delta": _profile_int(generated_page, "text_box_count")
        - _profile_int(ideal_page, "text_box_count"),
        "largest_picture_area_ratio_delta": round(
            _profile_number(generated_page, "largest_picture_area_ratio")
            - _profile_number(ideal_page, "largest_picture_area_ratio"),
            6,
        ),
        "picture_coverage_ratio_delta": round(
            _profile_number(generated_page, "picture_coverage_ratio")
            - _profile_number(ideal_page, "picture_coverage_ratio"),
            6,
        ),
        "top_geometry_mismatches": _top_geometry_mismatches(generated_page, ideal_page),
    }


def compare_pptx_structure(generated_path: Path, ideal_path: Path) -> dict:
    generated = Presentation(generated_path)
    ideal = Presentation(ideal_path)
    generated_profiles = _profile_by_page(profile_pptx_strategy(generated_path))
    ideal_profiles = _profile_by_page(profile_pptx_strategy(ideal_path))
    page_count = max(len(generated.slides), len(ideal.slides))
    pages: list[dict] = []
    for index in range(page_count):
        generated_counts = (
            _slide_counts(generated.slides[index])
            if index < len(generated.slides)
            else {"text_runs": 0, "pictures": 0, "shapes": 0, "tables": 0}
        )
        ideal_counts = (
            _slide_counts(ideal.slides[index])
            if index < len(ideal.slides)
            else {"text_runs": 0, "pictures": 0, "shapes": 0, "tables": 0}
        )
        generated_profile = generated_profiles.get(index + 1, {})
        ideal_profile = ideal_profiles.get(index + 1, {})
        pages.append(
            {
                "page_number": index + 1,
                **_strategy_delta(generated_profile, ideal_profile),
                "generated_text_runs": generated_counts["text_runs"],
                "ideal_text_runs": ideal_counts["text_runs"],
                "generated_pictures": generated_counts["pictures"],
                "ideal_pictures": ideal_counts["pictures"],
                "generated_shapes": generated_counts["shapes"],
                "ideal_shapes": ideal_counts["shapes"],
                "generated_tables": generated_counts["tables"],
                "ideal_tables": ideal_counts["tables"],
            }
        )
    return {
        "generated_slide_count": len(generated.slides),
        "ideal_slide_count": len(ideal.slides),
        "slide_count_delta": len(generated.slides) - len(ideal.slides),
        "pages": pages,
    }
